# Leo Lake
# SSF Analytics
# May 23, 2024
# this will pull tables from the sinto analytics powerpoints and implement them into excel
# read the READ ME for rules


import os
from pptx import Presentation
from openpyxl import load_workbook



def customer_date_extraction(presentation):
# from the file path open the presentation and assign blank values to the customer and date
    title_text = []
    customer = ''
    date = ''
# for each shape in the title slide check if it is a text box
# if yes add it to a list of the text for the title slide
    for shape in presentation.slides[0].shapes:
        if shape.has_text_frame:
            title_text.append(shape.text_frame.text.split('\n'))
# searches though the title text and finds the customer name and the date of the report
    for text_group in title_text:
        for text in text_group:
            if text.find('Customer: ') != -1:
                customer = text[text.index(': ')+2:]
            if text.find('Date: ') != -1:
                date = text[text.index(': ')+2:]
    return customer, date

def table_extraction(presentation):
# extracts the table objects from the presentation, max slides set to 50 to quicken loading times
# most tables will be on only the first couple slides
# loops through the slides and adds the table objects to a list
    max_slides = 50
    tables = []
    for i in range(min(max_slides, len(presentation.slides))):
        slide = presentation.slides[i]
        for shape in slide.shapes:
            if shape.has_table:
                tables.append(shape.table)
    return tables


def table_proccessing(table_list):
# processes the tables from function above, skips over header and first collum (where to find data in presentation)
# returns a large list of rows from the table
    data=[]
    for table in table_list:
        for i, row in enumerate(table.rows):
            if i == 0:
                continue
            row_data = [cell.text.replace('\x0b', '') for j, cell in enumerate(row.cells) if j != 0]
            data.append(row_data)
    return data


def combine_cust_date_data(cust_date, data):
# combine the customer and date into the row of the data, cust_date is a tuple
    customer = cust_date[0]
    date = cust_date[1]
    data = [[customer, date] + row for row in data]
    filtered_data = []
# delete the empty rows
    for row in data:
        if not all(val == '' for val in row[2:4]) and row[2] != 'None at this time':
            filtered_data.append(row)
    return filtered_data


def write_to_excel(excel_path, data):
# write all the tables to the excel file see path below
    workbook = load_workbook(excel_path)
    worksheet = workbook['anomalies']
    for row_data in data:
        row_num = worksheet.max_row + 1  # Determine the next available row
        for i, value in enumerate(row_data, start=1):
            worksheet.cell(row=row_num, column=i, value=value)
    workbook.save(excel_path)


def single_presentation_to_excel(presentation_path):
# this is the whole process for writing a presentation to excel
    sheet_path = r"V:\Srv\Sinto Analytics\_Anomaly Database\anomaly database.xlsx"
    print(presentation_path)
    presentation = Presentation(presentation_path)
    print('presentation complete')
    cust_date = customer_date_extraction(presentation)
    print('cust_date', cust_date)
    table_list = table_extraction(presentation)
    print('table_list', table_list)
    data = table_proccessing(table_list)
    data = combine_cust_date_data(cust_date, data)
    print('data', data)
    write_to_excel(sheet_path, data)
    print('writing to excel')


def filter_powerpoint_files(root_directory):
    filtered_files = []
    def should_process_directory(dir_path):
        # Split the directory path into its components and check if any start with '_'
        parts = dir_path.split(os.path.sep)
        return all(not part.startswith('_') and part != 'Archive' for part in parts if part)

    for root, dirs, files in os.walk(root_directory):
        # Check if the directory is a "weekly" folder within a customer folder
        if os.path.basename(root) == "Weekly" and should_process_directory(os.path.dirname(root)):
            # Filter PowerPoint files in the "weekly" folder
            for filename in files:
                if filename.endswith(".pptm") or filename.endswith(".ppt"):
                    if "YEAR" not in filename and "- 2020" not in filename and "- 2021" not in filename and 'Lifetime' not in filename and '~' not in filename and '.xlsx' not in filename:
                        filtered_files.append(os.path.join(root, filename))
        elif should_process_directory(root):
            # Filter PowerPoint files directly within the customer folder
            for filename in files:
                if filename.endswith(".pptm") or filename.endswith(".ppt"):
                    if "YEAR" not in filename and "- 2020" not in filename and "- 2021" not in filename and 'Lifetime' not in filename and '~' not in filename and '.xlsx' not in filename:

                        filtered_files.append(os.path.join(root, filename))

    return filtered_files


def update_paths_list(file_list):
# read all the presentations that have been transferred to excel, they are stored as a file path in the path below
# if the path from file list has a new path not in the text file then we will add it to the text file
# and pass it through to be processed by the rest of the code
    completed_anomalies_file_path = r"V:\Srv\Sinto Analytics\_Anomaly Database\completed_anomalies.txt"
    anomalies_to_complete = []
    with open(completed_anomalies_file_path, 'r') as completed_anomalies:
        completed_anomalies_list = completed_anomalies.read().splitlines()
    with open(completed_anomalies_file_path, 'a') as completed_anomalies:
        for anomaly_file in file_list:
            if anomaly_file not in completed_anomalies_list:
                anomalies_to_complete.append(anomaly_file.strip())
                completed_anomalies.write(anomaly_file + '\n')

    return anomalies_to_complete


# collect the list of files in sinto analytics, filter and check for repeats then pass to processing
root_directory_for_files = r"V:\Srv\Sinto Analytics"
filtered_files = filter_powerpoint_files(root_directory_for_files)
filtered_files = update_paths_list(filtered_files)
for file in filtered_files:
    single_presentation_to_excel(file)